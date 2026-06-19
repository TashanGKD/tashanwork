#!/usr/bin/env python3
"""Audit PPTX editable text against slide_spec.json editable_text."""

from __future__ import annotations

import argparse
import json
import re
from collections import Counter
from pathlib import Path

from pptx import Presentation


def _norm(text: str) -> str:
    text = str(text or "")
    text = re.sub(r"\s+", "", text)
    return text.strip()


def _expected_texts(spec_path: Path) -> dict[str, list[str]]:
    spec = json.loads(spec_path.read_text(encoding="utf-8"))
    out: dict[str, list[str]] = {}
    for index, slide in enumerate(spec.get("slides") or [], start=1):
        sid = str(slide.get("slide_id") or f"{index:02d}")
        texts = []
        for item in slide.get("editable_text") or []:
            if isinstance(item, dict) and str(item.get("text") or "").strip():
                texts.append(str(item["text"]))
        out[sid] = texts
    return out


def _actual_texts(pptx_path: Path) -> list[list[str]]:
    prs = Presentation(str(pptx_path))
    slides: list[list[str]] = []
    for slide in prs.slides:
        texts = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = shape.text_frame.text.strip()
                if text:
                    texts.append(text)
        slides.append(texts)
    return slides


def audit(spec_path: Path, pptx_path: Path) -> dict:
    expected_by_id = _expected_texts(spec_path)
    actual_slides = _actual_texts(pptx_path)
    issues: list[dict] = []
    slide_reports = []

    for index, (sid, expected_texts) in enumerate(expected_by_id.items(), start=1):
        actual_texts = actual_slides[index - 1] if index - 1 < len(actual_slides) else []
        expected_norm = Counter(_norm(text) for text in expected_texts if _norm(text))
        actual_norm = Counter(_norm(text) for text in actual_texts if _norm(text))

        missing = []
        for text in expected_norm:
            if actual_norm[text] < expected_norm[text]:
                missing.append(text)
                issues.append({"level": "error", "slide_id": sid, "message": f"missing editable text: {text}"})

        extra = []
        for text in actual_norm:
            if expected_norm[text] < actual_norm[text]:
                extra.append(text)
                issues.append({"level": "error", "slide_id": sid, "message": f"extra editable text not in spec: {text}"})

        duplicates = []
        for text, count in actual_norm.items():
            if count > max(1, expected_norm.get(text, 0)):
                duplicates.append({"text": text, "count": count})

        slide_reports.append({
            "slide_id": sid,
            "expected_count": len(expected_texts),
            "actual_count": len(actual_texts),
            "missing": missing,
            "extra": extra,
            "duplicates": duplicates,
            "actual_samples": actual_texts[:8],
        })

    if len(actual_slides) != len(expected_by_id):
        issues.append({
            "level": "error",
            "message": f"slide count mismatch: spec={len(expected_by_id)} pptx={len(actual_slides)}",
        })

    return {
        "status": "fail" if issues else "pass",
        "spec": str(spec_path),
        "pptx": str(pptx_path),
        "slide_count_spec": len(expected_by_id),
        "slide_count_pptx": len(actual_slides),
        "issue_count": len(issues),
        "issues": issues,
        "slides": slide_reports,
    }


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("spec", type=Path)
    parser.add_argument("pptx", type=Path)
    parser.add_argument("--out", type=Path)
    args = parser.parse_args()

    report = audit(args.spec, args.pptx)
    text = json.dumps(report, ensure_ascii=False, indent=2)
    if args.out:
        args.out.parent.mkdir(parents=True, exist_ok=True)
        args.out.write_text(text, encoding="utf-8")
    print(text)
    if report["status"] == "fail":
        raise SystemExit(1)


if __name__ == "__main__":
    main()
