#!/usr/bin/env python3
"""Audit whether a PPTX has real editable text and layered visual objects."""

from __future__ import annotations

import argparse
import json
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def audit(path: Path) -> dict:
    prs = Presentation(str(path))
    slides = []
    total_pictures = 0
    total_textboxes = 0
    total_empty_textboxes = 0
    total_shapes = 0

    for slide_index, slide in enumerate(prs.slides, start=1):
        pictures = 0
        textboxes = 0
        empty_textboxes = 0
        shape_types: dict[str, int] = {}
        text_samples: list[str] = []
        for shape in slide.shapes:
            total_shapes += 1
            shape_type = str(shape.shape_type)
            shape_types[shape_type] = shape_types.get(shape_type, 0) + 1
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                pictures += 1
            if getattr(shape, "has_text_frame", False):
                text = shape.text_frame.text.strip()
                if text:
                    textboxes += 1
                    if len(text_samples) < 5:
                        text_samples.append(text[:80])
                else:
                    empty_textboxes += 1
        total_pictures += pictures
        total_textboxes += textboxes
        total_empty_textboxes += empty_textboxes
        slides.append({
            "slide_index": slide_index,
            "pictures": pictures,
            "textboxes": textboxes,
            "empty_textboxes": empty_textboxes,
            "shape_types": shape_types,
            "text_samples": text_samples,
        })

    flattened_like = bool(len(prs.slides) > 0 and total_textboxes == 0 and total_pictures >= len(prs.slides))
    return {
        "pptx": str(path),
        "slide_count": len(prs.slides),
        "total_shapes": total_shapes,
        "total_pictures": total_pictures,
        "total_textboxes": total_textboxes,
        "total_empty_textboxes": total_empty_textboxes,
        "flattened_like": flattened_like,
        "slides": slides,
    }


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("pptx", type=Path)
    parser.add_argument("--out", type=Path)
    parser.add_argument("--min-slides", type=int, default=1)
    parser.add_argument("--min-pictures", type=int, default=1)
    parser.add_argument("--min-textboxes", type=int, default=1)
    parser.add_argument("--fail-flattened", action="store_true")
    args = parser.parse_args()

    report = audit(args.pptx)
    issues = []
    if report["slide_count"] < args.min_slides:
        issues.append(f"slide_count below {args.min_slides}")
    if report["total_pictures"] < args.min_pictures:
        issues.append(f"total_pictures below {args.min_pictures}")
    if report["total_textboxes"] < args.min_textboxes:
        issues.append(f"total_textboxes below {args.min_textboxes}")
    if args.fail_flattened and report["flattened_like"]:
        issues.append("pptx looks flattened: pictures exist but no real text boxes")
    report["status"] = "fail" if issues else "pass"
    report["issues"] = issues

    text = json.dumps(report, ensure_ascii=False, indent=2)
    if args.out:
        args.out.parent.mkdir(parents=True, exist_ok=True)
        args.out.write_text(text, encoding="utf-8")
    print(text)
    if report["status"] == "fail":
        raise SystemExit(1)


if __name__ == "__main__":
    main()
