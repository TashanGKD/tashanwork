#!/usr/bin/env python3
"""Self-test visual-deck-builder with layered editable deck fixtures."""

from __future__ import annotations

import argparse
import importlib.util
import json
import re
import shutil
import sys
from copy import deepcopy
from pathlib import Path

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

SCRIPT_DIR = Path(__file__).resolve().parent
SKILL_DIR = SCRIPT_DIR.parent
sys.path.insert(0, str(SCRIPT_DIR))

from compose_layered_deck import compose  # noqa: E402
from render_layered_preview import render as render_previews  # noqa: E402
from validate_visual_deck import validate  # noqa: E402
from audit_pptx_editability import audit as audit_pptx  # noqa: E402
from audit_pptx_text_against_spec import audit as audit_pptx_text  # noqa: E402
from audit_layered_layout import audit as audit_layout  # noqa: E402
from audit_visual_quality import audit as audit_visual_quality  # noqa: E402
from audit_visual_acceptance import audit as audit_visual_acceptance  # noqa: E402
from repair_layered_layout import repair as repair_layout  # noqa: E402

GORDEN_COMPOSE_PATH = SCRIPT_DIR / "gorden_image2pptx" / "compose_pptx.py"
_gorden_spec = importlib.util.spec_from_file_location("vdb_gorden_compose_pptx", GORDEN_COMPOSE_PATH)
if _gorden_spec is None or _gorden_spec.loader is None:
    raise RuntimeError(f"cannot load {GORDEN_COMPOSE_PATH}")
gorden_compose = importlib.util.module_from_spec(_gorden_spec)
_gorden_spec.loader.exec_module(gorden_compose)

ICON_COVERAGE_PATH = SCRIPT_DIR / "gorden_image2pptx" / "icon_coverage_audit.py"
_icon_coverage_spec = importlib.util.spec_from_file_location("vdb_icon_coverage_audit", ICON_COVERAGE_PATH)
if _icon_coverage_spec is None or _icon_coverage_spec.loader is None:
    raise RuntimeError(f"cannot load {ICON_COVERAGE_PATH}")
icon_coverage_audit = importlib.util.module_from_spec(_icon_coverage_spec)
_icon_coverage_spec.loader.exec_module(icon_coverage_audit)

FRAME_RESIDUE_PATH = SCRIPT_DIR / "gorden_image2pptx" / "frame_residue_audit.py"
_frame_residue_spec = importlib.util.spec_from_file_location("vdb_frame_residue_audit", FRAME_RESIDUE_PATH)
if _frame_residue_spec is None or _frame_residue_spec.loader is None:
    raise RuntimeError(f"cannot load {FRAME_RESIDUE_PATH}")
frame_residue_audit = importlib.util.module_from_spec(_frame_residue_spec)
_frame_residue_spec.loader.exec_module(frame_residue_audit)

BUILD_FRAME_RESIDUE_CONTRACT_PATH = SCRIPT_DIR / "gorden_image2pptx" / "build_frame_residue_contract.py"
_build_frame_residue_contract_spec = importlib.util.spec_from_file_location(
    "vdb_build_frame_residue_contract",
    BUILD_FRAME_RESIDUE_CONTRACT_PATH,
)
if _build_frame_residue_contract_spec is None or _build_frame_residue_contract_spec.loader is None:
    raise RuntimeError(f"cannot load {BUILD_FRAME_RESIDUE_CONTRACT_PATH}")
build_frame_residue_contract = importlib.util.module_from_spec(_build_frame_residue_contract_spec)
_build_frame_residue_contract_spec.loader.exec_module(build_frame_residue_contract)

PROMPT_PACK_PATH = SCRIPT_DIR / "gorden_image2pptx" / "build_extraction_prompt_pack.py"
_prompt_pack_spec = importlib.util.spec_from_file_location("vdb_build_extraction_prompt_pack", PROMPT_PACK_PATH)
if _prompt_pack_spec is None or _prompt_pack_spec.loader is None:
    raise RuntimeError(f"cannot load {PROMPT_PACK_PATH}")
build_extraction_prompt_pack = importlib.util.module_from_spec(_prompt_pack_spec)
_prompt_pack_spec.loader.exec_module(build_extraction_prompt_pack)


def _make_background(path: Path, color: tuple[int, int, int]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    im = Image.new("RGB", (1600, 900), color)
    draw = ImageDraw.Draw(im)
    draw.rectangle((0, 760, 1600, 900), fill=(230, 240, 238))
    draw.arc((-120, 600, 600, 1080), 195, 350, fill=(30, 120, 118), width=16)
    im.save(path)


def _make_frame(path: Path) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    im = Image.new("RGBA", (1600, 900), (0, 0, 0, 0))
    draw = ImageDraw.Draw(im)
    for i in range(3):
        x0 = 140 + i * 460
        draw.rounded_rectangle((x0, 250, x0 + 360, 650), radius=24, outline=(28, 96, 98, 255), width=8, fill=(255, 255, 255, 210))
        draw.rectangle((x0, 250, x0 + 360, 308), fill=(26, 120, 116, 235))
    draw.line((180, 710, 1420, 710), fill=(216, 157, 54, 255), width=5)
    im.save(path)


def _make_dense_frame(path: Path) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    im = Image.new("RGBA", (1600, 900), (0, 0, 0, 0))
    draw = ImageDraw.Draw(im)
    card_w = 330
    card_h = 130
    gap_x = 42
    start_x = 128
    start_y = 250
    for row in range(3):
        for col in range(4):
            x0 = start_x + col * (card_w + gap_x)
            y0 = start_y + row * 150
            draw.rounded_rectangle((x0, y0, x0 + card_w, y0 + card_h), radius=18, outline=(28, 96, 98, 255), width=6, fill=(255, 255, 255, 220))
            draw.rectangle((x0, y0, x0 + card_w, y0 + 38), fill=(38, 137, 130, 235))
    draw.line((180, 710, 1420, 710), fill=(216, 157, 54, 255), width=5)
    im.save(path)


def _make_chart_table_frame(path: Path) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    im = Image.new("RGBA", (1600, 900), (0, 0, 0, 0))
    draw = ImageDraw.Draw(im)
    draw.rounded_rectangle((80, 205, 965, 640), radius=22, outline=(28, 96, 98, 255), width=5, fill=(255, 255, 255, 225))
    draw.rectangle((120, 275, 900, 550), outline=(110, 135, 140, 180), width=2)
    for x in [250, 380, 510, 640, 770]:
        draw.line((x, 300, x, 550), fill=(205, 215, 215, 180), width=1)
    for y in [330, 385, 440, 495]:
        draw.line((120, y, 900, y), fill=(205, 215, 215, 180), width=1)
    bars = [(160, 470, 220, 550), (290, 410, 350, 550), (420, 365, 480, 550), (550, 330, 610, 550), (680, 385, 740, 550), (810, 315, 870, 550)]
    colors = [(38, 137, 130, 235), (215, 157, 54, 235), (70, 95, 135, 235), (38, 137, 130, 235), (215, 157, 54, 235), (70, 95, 135, 235)]
    for box, color in zip(bars, colors):
        draw.rounded_rectangle(box, radius=8, fill=color)
    draw.rounded_rectangle((1000, 205, 1500, 585), radius=22, outline=(28, 96, 98, 255), width=5, fill=(255, 255, 255, 225))
    draw.rectangle((1000, 205, 1500, 255), fill=(38, 137, 130, 235))
    for row in range(1, 6):
        y = 255 + row * 55
        draw.line((1025, y, 1475, y), fill=(205, 215, 215, 210), width=2)
    for x in [1160, 1300, 1400]:
        draw.line((x, 255, x, 585), fill=(205, 215, 215, 210), width=2)
    draw.rounded_rectangle((80, 670, 1500, 785), radius=18, outline=(28, 96, 98, 220), width=4, fill=(255, 255, 255, 210))
    for x, color in [(135, (38, 137, 130, 255)), (430, (215, 157, 54, 255)), (725, (70, 95, 135, 255))]:
        draw.rectangle((x, 708, x + 34, 742), fill=color)
    im.save(path)


def _make_icon(path: Path, color: tuple[int, int, int]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    im = Image.new("RGBA", (256, 256), (0, 0, 0, 0))
    draw = ImageDraw.Draw(im)
    draw.ellipse((32, 32, 224, 224), fill=color + (255,))
    draw.line((82, 132, 120, 170, 178, 82), fill=(255, 255, 255, 255), width=18)
    im.save(path)


def _make_plain_bg(path: Path) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    im = Image.new("RGB", (1600, 900), (248, 251, 250))
    draw = ImageDraw.Draw(im)
    draw.rounded_rectangle((80, 90, 1520, 810), radius=28, outline=(32, 111, 110), width=5, fill=(255, 255, 255))
    draw.rectangle((80, 90, 1520, 165), fill=(229, 243, 241))
    im.save(path)


def _dark_fraction(path: Path, box: tuple[int, int, int, int]) -> float:
    im = Image.open(path).convert("RGB").crop(box)
    width, height = im.size
    total = width * height
    if not total:
        return 0.0
    pixels = im.load()
    dark = 0
    for y in range(height):
        for x in range(width):
            r, g, b = pixels[x, y]
            if r < 45 and g < 45 and b < 45:
                dark += 1
    return dark / total


def _base_spec() -> dict:
    return {
        "deck": {
            "title": "AI科研助理方案",
            "language": "zh",
            "aspect_ratio": "16:9",
            "audience": "高校科研团队",
            "route": ["topic_only"],
            "style_brief": "clean executive visual report, layered editable PPTX",
        },
        "slides": [
            {
                "slide_id": "01",
                "purpose": "建立主题和核心判断",
                "title": "AI科研助理方案",
                "body_text": ["从检索到讲解形成可审计工作流"],
                "must_show": ["可审计", "科研助理"],
                "visual_brief": "layered editable slide with three-card framework",
                "editable_strategy": "layered_editable",
                "layer_prompts": {
                    "visual_target": "Generate the full-slide visual target with all intended text, rich hierarchy, and complete composition. This is the aesthetic reference for extraction, not the final editable PPT layer.",
                    "background": "Generate text-free and frame-free background only.",
                    "frame": "Using the visual target as edit target, extract structural cards, title bars, connectors, and chart scaffolding only. No ordinary text, no icons.",
                    "icons": "Using the visual target as edit target, extract movable icon and decoration elements only. No ordinary text.",
                },
                "editable_text": [
                    {
                        "text": "AI科研助理方案",
                        "role": "title",
                        "box": {"x": 0.07, "y": 0.08, "w": 0.58, "h": 0.10},
                        "style": {"font_size": 30, "bold": True, "color": "111111"},
                    }
                ],
                "evidence": [],
                "status": "layers_rendered",
            },
            {
                "slide_id": "02",
                "purpose": "说明执行路径",
                "title": "生成路径",
                "body_text": ["结构化规格、分层生成、PPTX合成、预览QA"],
                "must_show": ["slide_spec", "layered_deck", "layer_manifest"],
                "visual_brief": "layered editable pipeline slide",
                "editable_strategy": "layered_editable",
                "layer_prompts": {
                    "visual_target": "Generate the full-slide visual target with all intended text, rich hierarchy, and complete composition. This is the aesthetic reference for extraction, not the final editable PPT layer.",
                    "background": "Generate text-free and frame-free background only.",
                    "frame": "Using the visual target as edit target, extract pipeline cards, arrows, and QA matrix structure only. No text, no icons.",
                    "icons": "Using the visual target as edit target, extract small workflow icons only. No ordinary text.",
                },
                "editable_text": [
                    {
                        "text": "生成路径",
                        "role": "title",
                        "box": {"x": 0.07, "y": 0.08, "w": 0.45, "h": 0.10},
                        "style": {"font_size": 28, "bold": True, "color": "111111"},
                    }
                ],
                "evidence": [],
                "status": "layers_rendered",
            },
        ],
    }


def _text_from_editable(item: dict) -> dict:
    box = item["box"]
    style = item.get("style") or {}
    return {
        "text": item["text"],
        "x": box["x"],
        "y": box["y"],
        "w": box["w"],
        "h": box["h"],
        "size": style.get("font_size", 20),
        "bold": style.get("bold", False),
        "color": style.get("color", "111111"),
        "font": style.get("font_face", "Microsoft YaHei"),
        "role": item.get("role", ""),
    }


def _layered_deck_from_spec(spec: dict) -> dict:
    slides = []
    colors = [(42, 80, 102), (76, 91, 122)]
    for index, slide in enumerate(spec["slides"], start=1):
        sid = slide["slide_id"]
        slides.append({
            "slide_id": sid,
            "background": f"layers/{sid}/background.png",
            "frame": f"layers/{sid}/frame.png",
            "icons": [
                {"file": f"layers/{sid}/icons/check.png", "x": 0.74, "y": 0.42, "w": 0.08, "h": 0.14}
            ],
            "texts": [_text_from_editable(item) for item in slide["editable_text"]],
        })
    return {
        "slide_width_in": 13.333333,
        "slide_height_in": 7.5,
        "units": "fraction",
        "assets_dir": ".",
        "slides": slides,
    }


def _layer_manifest(spec: dict, backend: str = "codex-imagegen", status: str = "completed") -> dict:
    manifest_slides = []
    for slide in spec["slides"]:
        sid = slide["slide_id"]
        manifest_slides.append({
            "slide_id": sid,
            "layers": [
                {
                    "layer": "visual_target",
                    "backend": backend,
                    "prompt_file": f"prompts/{sid}-visual-target.md",
                    "generated_source": "self-test image-model output",
                    "copied_to": f"visual_targets/{sid}.png",
                    "status": status,
                },
                {
                    "layer": "background",
                    "backend": backend,
                    "prompt_file": f"prompts/{sid}-background.md",
                    "generated_source": "self-test image-model output",
                    "copied_to": f"layers/{sid}/background.png",
                    "status": status,
                },
                {
                    "layer": "frame",
                    "backend": backend,
                    "prompt_file": f"prompts/{sid}-frame.md",
                    "generated_source": "self-test image-model output",
                    "copied_to": f"layers/{sid}/frame.png",
                    "status": status,
                },
                {
                    "layer": "icons",
                    "backend": backend,
                    "prompt_file": f"prompts/{sid}-icons.md",
                    "generated_source": "self-test image-model output",
                    "copied_to": f"layers/{sid}/icons/check.png",
                    "status": status,
                },
            ],
        })
    return {"slides": manifest_slides}


def _dense_spec() -> dict:
    spec = _base_spec()
    spec["deck"]["title"] = "密集型可编辑方案页"
    spec["slides"] = [spec["slides"][0]]
    slide = spec["slides"][0]
    slide["slide_id"] = "01"
    slide["title"] = "密集型能力矩阵"
    slide["body_text"] = ["检索", "分析", "生成", "校验", "交付"]
    slide["must_show"] = ["证据链", "结构化", "可编辑", "预览QA", "发布"]
    slide["visual_brief"] = "dense editable dashboard with many cards, metrics, and icon markers"
    slide["editable_text"] = [
        {"text": "密集型能力矩阵", "role": "title", "box": {"x": 0.05, "y": 0.06, "w": 0.46, "h": 0.08}, "style": {"font_size": 28, "bold": True, "color": "111111"}},
        {"text": "面向复杂材料的分层可编辑PPT生成", "role": "subtitle", "box": {"x": 0.05, "y": 0.15, "w": 0.52, "h": 0.05}, "style": {"font_size": 14, "bold": False, "color": "245f5d"}},
    ]
    labels = ["需求解析", "资料抽取", "证据绑定", "背景层", "框架层", "图标层", "文本层", "PPTX合成", "预览QA", "泄漏扫描", "复杂版式", "发布门禁"]
    for idx, label in enumerate(labels):
        col = idx % 4
        row = idx // 4
        slide["editable_text"].append({
            "text": label,
            "role": "card_label",
            "box": {"x": 0.08 + col * 0.22, "y": 0.30 + row * 0.16, "w": 0.15, "h": 0.04},
            "style": {"font_size": 12, "bold": True, "color": "111111"},
        })
        slide["editable_text"].append({
            "text": f"检查项 {idx + 1:02d}",
            "role": "small_note",
            "box": {"x": 0.08 + col * 0.22, "y": 0.35 + row * 0.16, "w": 0.13, "h": 0.035},
            "style": {"font_size": 9, "bold": False, "color": "4b5563"},
        })
    return spec


def _dense_layered_deck(spec: dict) -> dict:
    texts = [_text_from_editable(item) for item in spec["slides"][0]["editable_text"]]
    icons = []
    for idx in range(12):
        col = idx % 4
        row = idx // 4
        icons.append({
            "file": "layers/01/icons/check.png",
            "x": 0.195 + col * 0.22,
            "y": 0.342 + row * 0.16,
            "w": 0.028,
            "h": 0.05,
        })
    return {
        "slide_width_in": 13.333333,
        "slide_height_in": 7.5,
        "units": "fraction",
        "assets_dir": ".",
        "slides": [{
            "slide_id": "01",
            "background": "layers/01/background.png",
            "frame": "layers/01/frame.png",
            "icons": icons,
            "texts": texts,
        }],
    }


def _chart_table_spec() -> dict:
    spec = _base_spec()
    spec["deck"]["title"] = "图表表格密集审计页"
    spec["slides"] = [spec["slides"][0]]
    slide = spec["slides"][0]
    slide["slide_id"] = "01"
    slide["title"] = "生成质量仪表盘"
    slide["body_text"] = ["同时覆盖趋势图、表格、legend 和指标说明"]
    slide["must_show"] = ["柱状图", "结果表格", "legend", "风险注释"]
    slide["visual_brief"] = "dense editable dashboard with chart, table, legend, and compact annotations"
    slide["editable_text"] = [
        {"text": "生成质量仪表盘", "role": "title", "box": {"x": 0.05, "y": 0.055, "w": 0.42, "h": 0.075}, "style": {"font_size": 26, "bold": True, "color": "111111"}},
        {"text": "复杂页可编辑性、证据链与视觉稳定性联检", "role": "subtitle", "box": {"x": 0.05, "y": 0.14, "w": 0.58, "h": 0.045}, "style": {"font_size": 13, "bold": False, "color": "245f5d"}},
        {"text": "月度通过率", "role": "chart_title", "box": {"x": 0.085, "y": 0.235, "w": 0.22, "h": 0.04}, "style": {"font_size": 13, "bold": True, "color": "111111"}},
        {"text": "发布门禁表", "role": "table_title", "box": {"x": 0.645, "y": 0.235, "w": 0.22, "h": 0.04}, "style": {"font_size": 13, "bold": True, "color": "ffffff"}},
    ]
    months = ["1月", "2月", "3月", "4月", "5月", "6月"]
    values = ["72", "81", "88", "93", "86", "96"]
    for idx, (month, value) in enumerate(zip(months, values)):
        x = 0.105 + idx * 0.081
        slide["editable_text"].append({"text": month, "role": "chart_axis_label", "box": {"x": x, "y": 0.62, "w": 0.045, "h": 0.03}, "style": {"font_size": 8, "color": "4b5563"}})
        slide["editable_text"].append({"text": value, "role": "chart_value_label", "box": {"x": x, "y": 0.31 + (100 - int(value)) * 0.002, "w": 0.045, "h": 0.03}, "style": {"font_size": 8, "bold": True, "color": "111111"}})
    headers = ["模块", "结果", "风险", "动作"]
    table_x = [0.645, 0.725, 0.805, 0.875]
    for idx, header in enumerate(headers):
        slide["editable_text"].append({"text": header, "role": "table_header", "box": {"x": table_x[idx], "y": 0.292, "w": 0.055, "h": 0.03}, "style": {"font_size": 8, "bold": True, "color": "245f5d"}})
    rows = [
        ["文本", "通过", "低", "保留"],
        ["图层", "通过", "中", "复核"],
        ["证据", "通过", "低", "归档"],
        ["样式", "警告", "中", "重排"],
        ["泄漏", "通过", "低", "发布"],
    ]
    for r, row in enumerate(rows):
        for c, cell in enumerate(row):
            slide["editable_text"].append({"text": cell, "role": "table_cell", "box": {"x": table_x[c], "y": 0.355 + r * 0.061, "w": 0.055, "h": 0.028}, "style": {"font_size": 8, "color": "111111"}})
    legend_positions = [(0.11, "结构通过"), (0.295, "需要复核"), (0.45, "高风险阻断")]
    for x, label in legend_positions:
        slide["editable_text"].append({"text": label, "role": "legend_label", "box": {"x": x, "y": 0.79, "w": 0.11, "h": 0.035}, "style": {"font_size": 9, "color": "111111"}})
    slide["editable_text"].append({"text": "说明：所有普通文字必须来自 slide_spec 并在 PPTX 中保持可编辑。", "role": "annotation", "box": {"x": 0.61, "y": 0.79, "w": 0.32, "h": 0.04}, "style": {"font_size": 9, "color": "4b5563"}})
    return spec


def _chart_table_layered_deck(spec: dict) -> dict:
    texts = [_text_from_editable(item) for item in spec["slides"][0]["editable_text"]]
    icons = [
        {"file": "layers/01/icons/check.png", "x": 0.085, "y": 0.785, "w": 0.02, "h": 0.036},
        {"file": "layers/01/icons/check.png", "x": 0.27, "y": 0.785, "w": 0.02, "h": 0.036},
        {"file": "layers/01/icons/check.png", "x": 0.425, "y": 0.785, "w": 0.02, "h": 0.036},
    ]
    return {
        "slide_width_in": 13.333333,
        "slide_height_in": 7.5,
        "units": "fraction",
        "assets_dir": ".",
        "slides": [{
            "slide_id": "01",
            "background": "layers/01/background.png",
            "frame": "layers/01/frame.png",
            "icons": icons,
            "texts": texts,
        }],
    }


def _write_case(root: Path, name: str, spec: dict, manifest: dict | None = None, layered: dict | None = None) -> tuple[Path, Path, Path]:
    case_dir = root / name
    case_dir.mkdir(parents=True, exist_ok=True)
    colors = [(242, 246, 244), (246, 244, 238)]
    for index, slide in enumerate(spec["slides"], start=1):
        sid = slide["slide_id"]
        _make_background(case_dir / "visual_targets" / f"{sid}.png", (236, 241, 245))
        _make_background(case_dir / "layers" / sid / "background.png", colors[(index - 1) % len(colors)])
        _make_frame(case_dir / "layers" / sid / "frame.png")
        _make_icon(case_dir / "layers" / sid / "icons" / "check.png", (28, 130, 116))
    spec_path = case_dir / "slide_spec.json"
    layered_path = case_dir / "layered_deck.json"
    manifest_path = case_dir / "layer_manifest.json"
    spec_path.write_text(json.dumps(spec, ensure_ascii=False, indent=2), encoding="utf-8")
    layered_path.write_text(json.dumps(layered or _layered_deck_from_spec(spec), ensure_ascii=False, indent=2), encoding="utf-8")
    manifest_path.write_text(json.dumps(manifest or _layer_manifest(spec), ensure_ascii=False, indent=2), encoding="utf-8")
    return spec_path, layered_path, manifest_path


def _expect_status(name: str, report: dict, expected: str) -> dict:
    return {
        "name": name,
        "ok": report["status"] == expected,
        "expected": expected,
        "actual": report["status"],
        "issues": report["issues"],
    }


def _inspect_pptx(path: Path) -> dict:
    prs = Presentation(path)
    pictures = 0
    textboxes = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                pictures += 1
            if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip():
                textboxes += 1
    return {"slides": len(prs.slides), "pictures": pictures, "textboxes": textboxes}


def _write_flattened_pptx(image_path: Path, out_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = int(round(13.333333 * 914400))
    prs.slide_height = int(round(7.5 * 914400))
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(image_path), 0, 0, width=prs.slide_width, height=prs.slide_height)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))


def _write_missing_text_pptx(image_path: Path, out_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = int(round(13.333333 * 914400))
    prs.slide_height = int(round(7.5 * 914400))
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(image_path), 0, 0, width=prs.slide_width, height=prs.slide_height)
    box = slide.shapes.add_textbox(0, 0, int(prs.slide_width / 2), int(prs.slide_height / 8))
    box.text_frame.text = "密集型能力矩阵"
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))


def _scan_static_contract() -> list[dict]:
    issues: list[dict] = []
    secret_pattern = re.compile(r"eyJ[A-Za-z0-9_-]+\.[A-Za-z0-9_-]+\.[A-Za-z0-9_-]+")
    for path in SKILL_DIR.rglob("*"):
        if not path.is_file() or path.suffix.lower() not in {".md", ".py", ".yaml", ".yml", ".json"}:
            continue
        if path.name == "self_test_visual_deck.py":
            continue
        text = path.read_text(encoding="utf-8", errors="ignore")
        if secret_pattern.search(text):
            issues.append({"file": str(path.relative_to(SKILL_DIR)), "message": "JWT-like secret found"})
    return issues


def run(workdir: Path) -> dict:
    if workdir.exists():
        shutil.rmtree(workdir)
    workdir.mkdir(parents=True)
    results: list[dict] = []

    spec_path, layered_path, manifest_path = _write_case(workdir, "valid-layered", _base_spec())
    results.append(_expect_status("valid layered spec validates", validate(spec_path, manifest_path, layered_path), "pass"))
    pptx_path = spec_path.parent / "out" / "deck-reconstructed-editable.pptx"
    compose(layered_path, pptx_path)
    preview_paths = render_previews(layered_path, spec_path.parent / "previews", width=1200, show_boxes=True)
    pptx_info = _inspect_pptx(pptx_path)
    results.append({
        "name": "layered PPTX structure",
        "ok": pptx_info["slides"] == 2 and pptx_info["pictures"] == 6 and pptx_info["textboxes"] == 2,
        "details": pptx_info,
    })
    results.append({
        "name": "layered preview renders",
        "ok": len(preview_paths) == 2 and all(path.exists() and path.stat().st_size > 1000 for path in preview_paths),
        "details": {"previews": [str(path) for path in preview_paths]},
    })

    dense_spec = _dense_spec()
    dense_layered = _dense_layered_deck(dense_spec)
    spec_path, layered_path, manifest_path = _write_case(workdir, "dense-layered-dashboard", dense_spec, layered=_dense_layered_deck(dense_spec))
    _make_dense_frame(spec_path.parent / "layers" / "01" / "frame.png")
    results.append(_expect_status("dense layered dashboard validates", validate(spec_path, manifest_path, layered_path), "pass"))
    dense_pptx = spec_path.parent / "out" / "dense-reconstructed-editable.pptx"
    compose(layered_path, dense_pptx)
    dense_preview_paths = render_previews(layered_path, spec_path.parent / "previews", width=1600, show_boxes=True)
    dense_info = _inspect_pptx(dense_pptx)
    dense_audit = audit_pptx(dense_pptx)
    dense_text_audit = audit_pptx_text(spec_path, dense_pptx)
    dense_layout_audit = audit_layout(layered_path)
    dense_visual_quality = audit_visual_quality(layered_path)
    results.append({
        "name": "dense layered PPTX structure",
        "ok": dense_info["slides"] == 1 and dense_info["pictures"] == 14 and dense_info["textboxes"] == len(dense_spec["slides"][0]["editable_text"]),
        "details": dense_info,
    })
    results.append({
        "name": "dense editable audit passes",
        "ok": (
            dense_audit["slide_count"] == 1
            and dense_audit["total_pictures"] == 14
            and dense_audit["total_textboxes"] == len(dense_spec["slides"][0]["editable_text"])
            and not dense_audit["flattened_like"]
        ),
        "details": dense_audit,
    })
    results.append({
        "name": "dense PPTX text matches spec",
        "ok": dense_text_audit["status"] == "pass" and dense_text_audit["issue_count"] == 0,
        "details": dense_text_audit,
    })
    results.append({
        "name": "dense layered layout audit warns on crowding",
        "ok": dense_layout_audit["status"] == "warn" and dense_layout_audit["issue_count"] > 0,
        "details": dense_layout_audit,
    })
    results.append({
        "name": "dense card-only visual quality audit warns",
        "ok": dense_visual_quality["status"] == "warn" and dense_visual_quality["issue_count"] > 0,
        "details": dense_visual_quality,
    })
    repaired_layered_path = spec_path.parent / "layered_deck.repaired.json"
    repair_report = repair_layout(layered_path, repaired_layered_path)
    repaired_layout_audit = audit_layout(repaired_layered_path)
    repaired_pptx = spec_path.parent / "out" / "dense-reconstructed-editable-repaired.pptx"
    compose(repaired_layered_path, repaired_pptx)
    repaired_preview_paths = render_previews(repaired_layered_path, spec_path.parent / "previews-repaired", width=1600, show_boxes=True)
    repaired_text_audit = audit_pptx_text(spec_path, repaired_pptx)
    repaired_edit_audit = audit_pptx(repaired_pptx)
    results.append({
        "name": "dense layout repair clears crowding",
        "ok": repair_report["change_count"] == 12 and repaired_layout_audit["status"] == "pass",
        "details": {"repair": repair_report, "layout": repaired_layout_audit},
    })
    results.append({
        "name": "repaired dense PPTX remains editable and spec-aligned",
        "ok": (
            repaired_text_audit["status"] == "pass"
            and repaired_edit_audit["total_textboxes"] == len(dense_spec["slides"][0]["editable_text"])
            and repaired_edit_audit["total_pictures"] == 14
            and not repaired_edit_audit["flattened_like"]
        ),
        "details": {"text": repaired_text_audit, "editability": repaired_edit_audit},
    })
    results.append({
        "name": "repaired dense preview renders",
        "ok": len(repaired_preview_paths) == 1 and repaired_preview_paths[0].exists() and repaired_preview_paths[0].stat().st_size > 1000,
        "details": {"previews": [str(path) for path in repaired_preview_paths]},
    })
    results.append({
        "name": "dense layered preview renders",
        "ok": len(dense_preview_paths) == 1 and dense_preview_paths[0].exists() and dense_preview_paths[0].stat().st_size > 1000,
        "details": {"previews": [str(path) for path in dense_preview_paths]},
    })

    chart_spec = _chart_table_spec()
    chart_layered = _chart_table_layered_deck(chart_spec)
    chart_spec_path, chart_layered_path, chart_manifest_path = _write_case(workdir, "dense-chart-table-legend", chart_spec, layered=chart_layered)
    _make_chart_table_frame(chart_spec_path.parent / "layers" / "01" / "frame.png")
    chart_pptx = chart_spec_path.parent / "out" / "chart-table-reconstructed-editable.pptx"
    compose(chart_layered_path, chart_pptx)
    chart_previews = render_previews(chart_layered_path, chart_spec_path.parent / "previews", width=1600, show_boxes=True)
    chart_validate = validate(chart_spec_path, chart_manifest_path, chart_layered_path)
    chart_layout_audit = audit_layout(chart_layered_path)
    chart_visual_quality = audit_visual_quality(chart_layered_path)
    chart_text_audit = audit_pptx_text(chart_spec_path, chart_pptx)
    chart_edit_audit = audit_pptx(chart_pptx)
    results.append(_expect_status("dense chart/table/legend validates", chart_validate, "pass"))
    results.append({
        "name": "dense chart/table/legend layout passes",
        "ok": chart_layout_audit["status"] == "pass" and chart_layout_audit["issue_count"] == 0,
        "details": chart_layout_audit,
    })
    results.append({
        "name": "dense chart/table/legend visual quality passes",
        "ok": chart_visual_quality["status"] == "pass" and chart_visual_quality["slides"][0]["has_chart"] and chart_visual_quality["slides"][0]["has_table"] and chart_visual_quality["slides"][0]["has_legend"],
        "details": chart_visual_quality,
    })
    results.append({
        "name": "dense chart/table/legend PPTX editable and text-aligned",
        "ok": (
            chart_text_audit["status"] == "pass"
            and chart_edit_audit["total_textboxes"] == len(chart_spec["slides"][0]["editable_text"])
            and chart_edit_audit["total_pictures"] == 5
            and not chart_edit_audit["flattened_like"]
        ),
        "details": {"text": chart_text_audit, "editability": chart_edit_audit},
    })
    results.append({
        "name": "dense chart/table/legend preview renders",
        "ok": len(chart_previews) == 1 and chart_previews[0].exists() and chart_previews[0].stat().st_size > 1000,
        "details": {"previews": [str(path) for path in chart_previews]},
    })

    flattened_pptx = spec_path.parent / "out" / "flattened-image-only.pptx"
    _write_flattened_pptx(spec_path.parent / "visual_targets" / "01.png", flattened_pptx)
    flattened_audit = audit_pptx(flattened_pptx)
    results.append({
        "name": "flattened image-only PPTX audit fails",
        "ok": flattened_audit["flattened_like"] and flattened_audit["total_textboxes"] == 0,
        "details": flattened_audit,
    })

    missing_text_pptx = spec_path.parent / "out" / "missing-text.pptx"
    _write_missing_text_pptx(spec_path.parent / "visual_targets" / "01.png", missing_text_pptx)
    missing_text_audit = audit_pptx_text(spec_path, missing_text_pptx)
    results.append({
        "name": "PPTX text missing from spec audit fails",
        "ok": missing_text_audit["status"] == "fail" and missing_text_audit["issue_count"] > 0,
        "details": missing_text_audit,
    })

    overlapping_layered = _layered_deck_from_spec(_base_spec())
    overlapping_layered["slides"][0]["texts"].append(deepcopy(overlapping_layered["slides"][0]["texts"][0]))
    overlapping_layered["slides"][0]["texts"][-1]["text"] = "重叠文本"
    spec_path, layered_path, manifest_path = _write_case(workdir, "overlapping-text-layout", _base_spec(), layered=overlapping_layered)
    overlap_layout_audit = audit_layout(layered_path)
    results.append({
        "name": "overlapping text layout audit fails",
        "ok": overlap_layout_audit["status"] == "fail" and overlap_layout_audit["issue_count"] > 0,
        "details": overlap_layout_audit,
    })

    grounded = deepcopy(_base_spec())
    grounded["deck"]["route"] = ["source_grounded", "style_reference"]
    grounded["deck"]["reference_guard"] = {
        "style_reference_policy": "layout_density_only",
        "allowed_text_source": "slide_spec",
        "forbidden_terms": ["OpenAI", "GPT-4", "ChatGPT", "2023"],
    }
    for slide in grounded["slides"]:
        slide["evidence"] = [{"source_id": "doc-1", "locator": "p.1", "claim": slide["body_text"][0], "confidence": "direct"}]
    spec_path, layered_path, manifest_path = _write_case(workdir, "grounded-with-evidence", grounded)
    results.append(_expect_status("source grounded with evidence passes", validate(spec_path, manifest_path, layered_path), "pass"))

    style_no_guard = deepcopy(_base_spec())
    style_no_guard["deck"]["route"] = ["style_reference"]
    spec_path, layered_path, manifest_path = _write_case(workdir, "style-reference-no-guard", style_no_guard)
    results.append(_expect_status("style reference without guard fails", validate(spec_path, manifest_path, layered_path), "fail"))

    style_forbidden_content = deepcopy(_base_spec())
    style_forbidden_content["deck"]["route"] = ["style_reference"]
    style_forbidden_content["deck"]["reference_guard"] = {
        "style_reference_policy": "layout_density_only",
        "allowed_text_source": "slide_spec",
        "forbidden_terms": ["OpenAI"],
    }
    style_forbidden_content["slides"][0]["body_text"] = ["OpenAI 示例语义不应迁入"]
    spec_path, layered_path, manifest_path = _write_case(workdir, "style-reference-forbidden-content", style_forbidden_content)
    results.append(_expect_status("style reference forbidden content fails", validate(spec_path, manifest_path, layered_path), "fail"))

    existing_edit = deepcopy(_base_spec())
    existing_edit["deck"]["route"] = ["existing_slide_edit"]
    for slide in existing_edit["slides"]:
        slide["source_visual_target"] = f"visual_targets/{slide['slide_id']}.png"
    spec_path, layered_path, manifest_path = _write_case(workdir, "existing-slide-edit-with-target", existing_edit)
    results.append(_expect_status("existing slide edit with source target passes", validate(spec_path, manifest_path, layered_path), "pass"))

    existing_edit_missing = deepcopy(_base_spec())
    existing_edit_missing["deck"]["route"] = ["existing_slide_edit"]
    spec_path, layered_path, manifest_path = _write_case(workdir, "existing-slide-edit-missing-target", existing_edit_missing)
    results.append(_expect_status("existing slide edit missing source target fails", validate(spec_path, manifest_path, layered_path), "fail"))

    cjk_spec_path, cjk_layered_path, cjk_manifest_path = _write_case(workdir, "中文路径-复杂场景", _base_spec())
    cjk_pptx = cjk_spec_path.parent / "out" / "中文输出.pptx"
    compose(cjk_layered_path, cjk_pptx)
    cjk_info = _inspect_pptx(cjk_pptx)
    results.append({
        "name": "Chinese path layered compose passes",
        "ok": cjk_info["slides"] == 2 and cjk_info["pictures"] == 6 and cjk_info["textboxes"] == 2,
        "details": cjk_info,
    })

    gorden_dir = workdir / "gorden-compose-cjk-fit"
    _make_plain_bg(gorden_dir / "background.png")
    gorden_layout = {
        "slide_width_in": 13.333,
        "slide_height_in": 7.5,
        "units": "fraction",
        "assets_dir": str(gorden_dir),
        "slides": [
            {
                "background": "background.png",
                "texts": [
                    {
                        "text": "中文文本自适应预览回归",
                        "x": 0.08,
                        "y": 0.10,
                        "w": 0.58,
                        "h": 0.08,
                        "size": 34,
                        "bold": True,
                        "color": "#176D6A",
                        "fit_text": True,
                    },
                    {
                        "text": "这是一段用于检查密集页面中文字体、自动缩小和预览渲染的长文本，不能出现黑块或明显爆框。",
                        "x": 0.10,
                        "y": 0.30,
                        "w": 0.32,
                        "h": 0.16,
                        "size": 20,
                        "color": "#1F2C2F",
                        "fit_text": True,
                    },
                ],
            }
        ],
    }
    gorden_layout_path = gorden_dir / "layout.json"
    gorden_layout_path.write_text(json.dumps(gorden_layout, ensure_ascii=False, indent=2), encoding="utf-8")
    gorden_pptx = gorden_dir / "out" / "gorden-compose-cjk-fit.pptx"
    gorden_preview_dir = gorden_dir / "preview"
    gorden_deck = gorden_compose._load_deck(gorden_layout_path)
    gorden_compose.build_pptx(gorden_deck, gorden_pptx)
    gorden_compose.render_previews(gorden_deck, gorden_preview_dir)
    gorden_preview = gorden_preview_dir / "slide_01.png"
    gorden_info = audit_pptx(gorden_pptx)
    dark_fraction = _dark_fraction(gorden_preview, (120, 80, 1060, 430))
    results.append({
        "name": "Gorden compose CJK preview uses readable fitted text",
        "ok": (
            gorden_info["slide_count"] == 1
            and gorden_info["total_textboxes"] == 2
            and not gorden_info["flattened_like"]
            and gorden_preview.exists()
            and dark_fraction < 0.18
        ),
        "details": {
            "pptx": str(gorden_pptx),
            "preview": str(gorden_preview),
            "textboxes": gorden_info["total_textboxes"],
            "dark_fraction": round(dark_fraction, 4),
        },
    })

    icon_coverage_dir = workdir / "icon-coverage-audit"
    icon_coverage_dir.mkdir(parents=True, exist_ok=True)
    icon_layout = {
        "slide_width_in": 13.333,
        "slide_height_in": 7.5,
        "units": "fraction",
        "ref_width": 1600,
        "ref_height": 900,
        "slides": [
            {
                "icons": [
                    {"name": "top_target", "source_bbox": [100, 50, 40, 40]},
                    {"name": "top_layers", "source_bbox": [180, 50, 40, 40]},
                    {"name": "bottom_lock", "source_bbox": [1320, 800, 55, 55]},
                ],
                "texts": [],
            }
        ],
    }
    icon_expected_pass = {
        "overall_min_icons": 3,
        "regions": [
            {
                "name": "top_nav",
                "bbox": [80, 30, 180, 90],
                "min_icons": 2,
                "required_names": ["top_target", "top_layers"],
            },
            {
                "name": "bottom_status",
                "bbox": [1260, 760, 180, 120],
                "min_icons": 1,
                "required_names": [["bottom_*", "*lock"]],
            },
        ],
    }
    icon_expected_fail = deepcopy(icon_expected_pass)
    icon_expected_fail["regions"][1]["min_icons"] = 2
    icon_expected_fail["regions"][1]["required_names"] = ["bottom_lock", "bottom_flow"]
    icon_layout_path = icon_coverage_dir / "layout.json"
    icon_expected_pass_path = icon_coverage_dir / "expected-pass.json"
    icon_expected_fail_path = icon_coverage_dir / "expected-fail.json"
    icon_layout_path.write_text(json.dumps(icon_layout, ensure_ascii=False, indent=2), encoding="utf-8")
    icon_expected_pass_path.write_text(json.dumps(icon_expected_pass, ensure_ascii=False, indent=2), encoding="utf-8")
    icon_expected_fail_path.write_text(json.dumps(icon_expected_fail, ensure_ascii=False, indent=2), encoding="utf-8")
    icon_pass_report = icon_coverage_audit.audit(icon_layout_path, icon_expected_pass_path)
    icon_fail_report = icon_coverage_audit.audit(icon_layout_path, icon_expected_fail_path)
    results.append({
        "name": "icon coverage audit passes complete source regions",
        "ok": icon_pass_report["status"] == "pass" and icon_pass_report["issue_count"] == 0,
        "details": icon_pass_report,
    })
    results.append({
        "name": "icon coverage audit catches missing dense region icons",
        "ok": icon_fail_report["status"] == "fail" and icon_fail_report["issue_count"] >= 2,
        "details": icon_fail_report,
    })

    residue_dir = workdir / "frame-residue-audit"
    residue_dir.mkdir(parents=True, exist_ok=True)
    clean_frame = residue_dir / "frame-clean.png"
    dirty_frame = residue_dir / "frame-dirty.png"
    for path, dirty in [(clean_frame, False), (dirty_frame, True)]:
        im = Image.new("RGBA", (1600, 900), (0, 0, 0, 0))
        draw = ImageDraw.Draw(im)
        draw.rounded_rectangle((80, 740, 1520, 870), radius=18, outline=(28, 96, 98, 255), width=5, fill=(255, 255, 255, 210))
        if dirty:
            draw.ellipse((300, 755, 400, 855), fill=(0, 108, 104, 245))
        im.save(path)
    residue_regions = {
        "regions": [
            {
                "name": "bottom_status",
                "bbox": [60, 720, 1480, 170],
                "forbidden_residue": [
                    {
                        "name": "no_status_icon_bases_in_frame",
                        "detector": "saturated_components",
                        "color_family": "teal_green",
                        "min_area": 2500,
                        "min_side": 40,
                        "max_aspect": 3.0
                    }
                ],
            }
        ]
    }
    residue_regions_path = residue_dir / "regions.json"
    residue_regions_path.write_text(json.dumps(residue_regions, ensure_ascii=False, indent=2), encoding="utf-8")
    residue_no_checks = {
        "regions": [
            {
                "name": "bottom_bbox_without_forbidden_residue_contract",
                "bbox": [60, 720, 1480, 170],
            }
        ]
    }
    residue_no_checks_path = residue_dir / "regions-no-checks.json"
    residue_no_checks_path.write_text(json.dumps(residue_no_checks, ensure_ascii=False, indent=2), encoding="utf-8")
    residue_layout = {
        "slide_width_in": 13.333,
        "slide_height_in": 7.5,
        "units": "fraction",
        "ref_width": 1600,
        "ref_height": 900,
        "assets_dir": str(residue_dir),
        "slides": [{"background": "missing.png", "frame": "frame-clean.png", "icons": [], "texts": []}],
    }
    residue_layout_path = residue_dir / "layout-clean.json"
    residue_layout_path.write_text(json.dumps(residue_layout, ensure_ascii=False, indent=2), encoding="utf-8")
    clean_residue_report = frame_residue_audit.audit(residue_layout_path, residue_regions_path)
    dirty_layout = deepcopy(residue_layout)
    dirty_layout["slides"][0]["frame"] = "frame-dirty.png"
    dirty_layout_path = residue_dir / "layout-dirty.json"
    dirty_layout_path.write_text(json.dumps(dirty_layout, ensure_ascii=False, indent=2), encoding="utf-8")
    dirty_residue_report = frame_residue_audit.audit(dirty_layout_path, residue_regions_path)
    no_check_residue_report = frame_residue_audit.audit(dirty_layout_path, residue_no_checks_path)
    results.append({
        "name": "frame residue audit passes clean frame regions",
        "ok": clean_residue_report["status"] == "pass",
        "details": clean_residue_report,
    })
    results.append({
        "name": "frame residue audit catches icon bases in frame layer",
        "ok": dirty_residue_report["status"] == "fail" and dirty_residue_report["regions"][0]["suspicious_count"] == 1,
        "details": dirty_residue_report,
    })
    results.append({
        "name": "frame residue audit skips bbox-only regions without forbidden residue contract",
        "ok": (
            no_check_residue_report["status"] == "pass"
            and no_check_residue_report["regions"][0]["skipped"] is True
            and no_check_residue_report["regions"][0]["check_count"] == 0
        ),
        "details": no_check_residue_report,
    })
    skeleton_contract = build_frame_residue_contract.build_contract(icon_layout_path)
    explicit_contract = build_frame_residue_contract.build_contract(
        icon_layout_path,
        color_family=["teal_green"],
        padding=12,
    )
    results.append({
        "name": "frame residue contract builder defaults to bbox skeleton without hardcoded color checks",
        "ok": (
            skeleton_contract["mode"] == "bbox_skeleton_review_required"
            and skeleton_contract["regions"]
            and "forbidden_residue" not in skeleton_contract["regions"][0]
        ),
        "details": skeleton_contract,
    })
    results.append({
        "name": "frame residue contract builder adds explicit checks only when color family is declared",
        "ok": (
            explicit_contract["mode"] == "explicit_forbidden_residue"
            and explicit_contract["regions"]
            and explicit_contract["regions"][0]["forbidden_residue"][0]["color_family"] == "teal_green"
        ),
        "details": explicit_contract,
    })

    icon_fail_report_path = icon_coverage_dir / "icon-coverage-fail-report.json"
    icon_fail_report_path.write_text(json.dumps(icon_fail_report, ensure_ascii=False, indent=2), encoding="utf-8")
    prompt_pack_dir = icon_coverage_dir / "prompt-pack"
    prompt_pack = build_extraction_prompt_pack.build_pack(
        icon_fail_report_path,
        icon_expected_fail_path,
        prompt_pack_dir,
        deck_language="zh",
        chroma="#ff00ff",
    )
    prompt_files = [Path(item["prompt_file"]) for item in prompt_pack["prompts"]]
    prompt_text = "\n".join(path.read_text(encoding="utf-8") for path in prompt_files)
    results.append({
        "name": "coverage prompt pack targets failed icon/frame regions",
        "ok": (
            prompt_pack["status"] == "pass"
            and prompt_pack["prompt_count"] >= 2
            and any(item["kind"] == "icons" for item in prompt_pack["prompts"])
            and any(item["kind"] == "frame_detail" for item in prompt_pack["prompts"])
            and "only edit target" in prompt_text
            and "Do not infer from this text alone" in prompt_text
            and "bottom_flow" in prompt_text
        ),
        "details": prompt_pack,
    })

    planning = deepcopy(_base_spec())
    for slide in planning["slides"]:
        slide["status"] = "planned"
    spec_path, layered_path, manifest_path = _write_case(workdir, "planning", planning)
    results.append(_expect_status("planning spec allow-planning passes", validate(spec_path, None, None, allow_planning=True), "pass"))

    missing_deck = deepcopy(_base_spec())
    del missing_deck["deck"]["title"]
    spec_path, layered_path, manifest_path = _write_case(workdir, "missing-deck-title", missing_deck)
    results.append(_expect_status("missing deck title fails", validate(spec_path, manifest_path, layered_path), "fail"))

    invalid_strategy = deepcopy(_base_spec())
    invalid_strategy["slides"][0]["editable_strategy"] = "legacy_single_image"
    spec_path, layered_path, manifest_path = _write_case(workdir, "invalid-strategy", invalid_strategy)
    results.append(_expect_status("image-only strategy fails", validate(spec_path, manifest_path, layered_path), "fail"))

    missing_prompt = deepcopy(_base_spec())
    del missing_prompt["slides"][0]["layer_prompts"]["frame"]
    spec_path, layered_path, manifest_path = _write_case(workdir, "missing-layer-prompt", missing_prompt)
    results.append(_expect_status("missing frame prompt fails", validate(spec_path, manifest_path, layered_path), "fail"))

    missing_target_prompt = deepcopy(_base_spec())
    del missing_target_prompt["slides"][0]["layer_prompts"]["visual_target"]
    spec_path, layered_path, manifest_path = _write_case(workdir, "missing-visual-target-prompt", missing_target_prompt)
    results.append(_expect_status("missing visual target prompt fails", validate(spec_path, manifest_path, layered_path), "fail"))

    no_evidence = deepcopy(_base_spec())
    no_evidence["deck"]["route"] = ["source_grounded"]
    spec_path, layered_path, manifest_path = _write_case(workdir, "source-no-evidence", no_evidence)
    results.append(_expect_status("source grounded without evidence fails", validate(spec_path, manifest_path, layered_path), "fail"))

    duplicate = deepcopy(_base_spec())
    duplicate["slides"][1]["slide_id"] = duplicate["slides"][0]["slide_id"]
    spec_path, layered_path, manifest_path = _write_case(workdir, "duplicate-slide-id", duplicate)
    results.append(_expect_status("duplicate slide id fails", validate(spec_path, manifest_path, layered_path), "fail"))

    invalid_status = deepcopy(_base_spec())
    invalid_status["slides"][0]["status"] = "done"
    spec_path, layered_path, manifest_path = _write_case(workdir, "invalid-status", invalid_status)
    results.append(_expect_status("invalid status fails", validate(spec_path, manifest_path, layered_path), "fail"))

    bad_box = deepcopy(_base_spec())
    bad_box["slides"][0]["editable_text"][0]["box"] = {"x": 0.80, "y": 0.10, "w": 0.40, "h": 0.12}
    spec_path, layered_path, manifest_path = _write_case(workdir, "bad-editable-box", bad_box)
    results.append(_expect_status("bad editable text box fails", validate(spec_path, manifest_path, layered_path), "fail"))

    spec_path, layered_path, manifest_path = _write_case(workdir, "missing-layered-deck", _base_spec())
    results.append(_expect_status("missing layered_deck fails", validate(spec_path, manifest_path, spec_path.parent / "missing.json"), "fail"))

    spec_path, layered_path, manifest_path = _write_case(workdir, "missing-background-file", _base_spec())
    (spec_path.parent / "layers" / "01" / "background.png").unlink()
    results.append(_expect_status("missing background layer fails", validate(spec_path, manifest_path, layered_path), "fail"))

    spec_path, layered_path, manifest_path = _write_case(workdir, "corrupt-frame", _base_spec())
    (spec_path.parent / "layers" / "01" / "frame.png").write_text("not a png", encoding="utf-8")
    results.append(_expect_status("corrupt frame layer fails", validate(spec_path, manifest_path, layered_path), "fail"))

    spec_path, layered_path, manifest_path = _write_case(workdir, "invalid-backend", _base_spec(), _layer_manifest(_base_spec(), backend="PIL"))
    results.append(_expect_status("non-image layer backend fails", validate(spec_path, manifest_path, layered_path), "fail"))

    spec_path, layered_path, manifest_path = _write_case(workdir, "blocked-layer", _base_spec(), _layer_manifest(_base_spec(), status="blocked"))
    results.append(_expect_status("blocked layer fails", validate(spec_path, manifest_path, layered_path), "fail"))

    missing_layer_manifest = _layer_manifest(_base_spec())
    missing_layer_manifest["slides"][0]["layers"] = [
        layer for layer in missing_layer_manifest["slides"][0]["layers"] if layer["layer"] != "background"
    ]
    spec_path, layered_path, manifest_path = _write_case(workdir, "manifest-missing-background", _base_spec(), missing_layer_manifest)
    results.append(_expect_status("layer manifest missing background fails", validate(spec_path, manifest_path, layered_path), "fail"))

    missing_target_manifest = _layer_manifest(_base_spec())
    missing_target_manifest["slides"][0]["layers"] = [
        layer for layer in missing_target_manifest["slides"][0]["layers"] if layer["layer"] != "visual_target"
    ]
    spec_path, layered_path, manifest_path = _write_case(workdir, "manifest-missing-visual-target", _base_spec(), missing_target_manifest)
    results.append(_expect_status("layer manifest missing visual target fails", validate(spec_path, manifest_path, layered_path), "fail"))

    layered_no_text = _layered_deck_from_spec(_base_spec())
    layered_no_text["slides"][0]["texts"] = []
    spec_path, layered_path, manifest_path = _write_case(workdir, "layered-no-text", _base_spec(), layered=_layered_deck_from_spec(_base_spec()))
    layered_path.write_text(json.dumps(layered_no_text, ensure_ascii=False, indent=2), encoding="utf-8")
    results.append(_expect_status("layered slide without real text fails", validate(spec_path, manifest_path, layered_path), "fail"))

    layered_bad_icon = _layered_deck_from_spec(_base_spec())
    layered_bad_icon["slides"][0]["icons"][0]["x"] = 0.98
    spec_path, layered_path, manifest_path = _write_case(workdir, "bad-icon-box", _base_spec(), layered=layered_bad_icon)
    results.append(_expect_status("bad icon placement fails", validate(spec_path, manifest_path, layered_path), "fail"))
    try:
        compose(layered_path, spec_path.parent / "out" / "bad-icon.pptx")
        composer_rejected_bad_icon = False
    except SystemExit:
        composer_rejected_bad_icon = True
    results.append({"name": "bad icon placement fails compose", "ok": composer_rejected_bad_icon})

    unsupported_ratio = deepcopy(_base_spec())
    unsupported_ratio["deck"]["aspect_ratio"] = "21:9"
    spec_path, layered_path, manifest_path = _write_case(workdir, "unsupported-ratio", unsupported_ratio)
    results.append(_expect_status("unsupported aspect ratio warns", validate(spec_path, manifest_path, layered_path), "warn"))

    placeholder = deepcopy(_base_spec())
    placeholder["slides"][0]["body_text"] = ["TODO placeholder text for warning path"]
    spec_path, layered_path, manifest_path = _write_case(workdir, "placeholder-warning", placeholder)
    results.append(_expect_status("placeholder terms warn", validate(spec_path, manifest_path, layered_path), "warn"))

    visual_acceptance_dir = workdir / "visual-acceptance"
    visual_acceptance_dir.mkdir(parents=True, exist_ok=True)
    visual_pass = {
        "slides": [
            {
                "slide_id": "01",
                "target_quality": "pass",
                "preview_readability": "pass",
                "text_rendering": "pass",
                "layer_fidelity": "pass",
                "icon_coverage": "pass",
                "semantic_drift": "pass",
                "overall": "pass",
                "notes": ["Target and reconstructed preview are visually aligned; editable Chinese text is readable."],
            }
        ]
    }
    visual_pass_path = visual_acceptance_dir / "visual-review-pass.json"
    visual_pass_path.write_text(json.dumps(visual_pass, ensure_ascii=False, indent=2), encoding="utf-8")
    visual_pass_report = audit_visual_acceptance(visual_pass_path)
    results.append({
        "name": "visual acceptance review passes clean deck",
        "ok": visual_pass_report["status"] == "pass" and visual_pass_report["issue_count"] == 0,
        "details": visual_pass_report,
    })

    visual_fail = {
        "slides": [
            {
                "slide_id": "01",
                "target_quality": "pass",
                "preview_readability": "fail",
                "text_rendering": "fail",
                "layer_fidelity": "fail",
                "icon_coverage": "fail",
                "semantic_drift": "warn",
                "overall": "fail",
                "notes": [
                    "Preview renderer shows CJK text as black blocks.",
                    "Several icons are missing or duplicated.",
                    "Frame layer does not match visual target closely enough for release.",
                ],
            }
        ]
    }
    visual_fail_compare = {
        "mean_abs_diff_0_255": 21.98,
        "changed_pixel_fraction_threshold_32": 0.146,
        "changed_pixel_fraction_threshold_64": 0.11,
    }
    visual_fail_path = visual_acceptance_dir / "visual-review-fail.json"
    visual_fail_compare_path = visual_acceptance_dir / "visual-compare-fail.json"
    visual_fail_path.write_text(json.dumps(visual_fail, ensure_ascii=False, indent=2), encoding="utf-8")
    visual_fail_compare_path.write_text(json.dumps(visual_fail_compare, ensure_ascii=False, indent=2), encoding="utf-8")
    visual_fail_report = audit_visual_acceptance(visual_fail_path, compare_path=visual_fail_compare_path)
    results.append({
        "name": "visual acceptance catches editable but ugly reconstruction",
        "ok": visual_fail_report["status"] == "fail" and visual_fail_report["issue_count"] >= 4,
        "details": visual_fail_report,
    })

    static_issues = _scan_static_contract()
    results.append({"name": "static secret scan", "ok": not static_issues, "issues": static_issues})

    passed = sum(1 for result in results if result["ok"])
    return {"status": "pass" if passed == len(results) else "fail", "passed": passed, "total": len(results), "results": results}


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--workdir", type=Path, default=Path("work/visual-deck-builder-self-test"))
    parser.add_argument("--out", type=Path)
    args = parser.parse_args()

    report = run(args.workdir)
    text = json.dumps(report, ensure_ascii=False, indent=2)
    if args.out:
        args.out.parent.mkdir(parents=True, exist_ok=True)
        args.out.write_text(text, encoding="utf-8")
    print(text)
    if report["status"] != "pass":
        raise SystemExit(1)


if __name__ == "__main__":
    main()
