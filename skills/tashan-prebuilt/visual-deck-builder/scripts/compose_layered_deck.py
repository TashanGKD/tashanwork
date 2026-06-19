#!/usr/bin/env python3
"""Compose a layered editable PPTX from visual-deck-builder layered_deck.json."""

from __future__ import annotations

import argparse
import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


def _hex_color(value: str) -> RGBColor:
    value = (value or "111111").strip().lstrip("#")
    if len(value) == 3:
        value = "".join(ch * 2 for ch in value)
    if len(value) != 6:
        value = "111111"
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def _resolve(base: Path, maybe_path: str | None) -> Path | None:
    if not maybe_path:
        return None
    p = Path(maybe_path)
    if not p.is_absolute():
        p = base / p
    return p


def _fraction(item: dict, key: str) -> float:
    try:
        value = float(item[key])
    except KeyError as exc:
        raise SystemExit(f"missing {key} in positioned layer item") from exc
    except (TypeError, ValueError) as exc:
        raise SystemExit(f"{key} must be numeric in positioned layer item") from exc
    if not 0 <= value <= 1:
        raise SystemExit(f"{key} outside [0, 1] in positioned layer item")
    return value


def _box(item: dict) -> tuple[float, float, float, float]:
    x = _fraction(item, "x")
    y = _fraction(item, "y")
    w = _fraction(item, "w")
    h = _fraction(item, "h")
    if w <= 0 or h <= 0:
        raise SystemExit("positioned layer width/height must be positive")
    if x + w > 1 or y + h > 1:
        raise SystemExit("positioned layer extends beyond slide canvas")
    return x, y, w, h


def _add_text(slide, item: dict, slide_w: float, slide_h: float) -> None:
    text = str(item.get("text") or "")
    if not text.strip():
        return
    x, y, w, h = _box(item)
    tx = slide.shapes.add_textbox(Inches(x * slide_w), Inches(y * slide_h), Inches(w * slide_w), Inches(h * slide_h))
    frame = tx.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Pt(float(item.get("margin_left", 0)))
    frame.margin_right = Pt(float(item.get("margin_right", 0)))
    frame.margin_top = Pt(float(item.get("margin_top", 0)))
    frame.margin_bottom = Pt(float(item.get("margin_bottom", 0)))
    valign = str(item.get("valign") or "top")
    frame.vertical_anchor = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "center": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}.get(valign, MSO_ANCHOR.TOP)

    para = frame.paragraphs[0]
    para.text = text
    para.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY}.get(str(item.get("align") or "left"), PP_ALIGN.LEFT)
    if item.get("line_spacing") is not None:
        para.line_spacing = float(item["line_spacing"])
    font = para.font
    font.size = Pt(float(item.get("size", item.get("font_size", 18))))
    font.bold = bool(item.get("bold", False))
    font.italic = bool(item.get("italic", False))
    font.color.rgb = _hex_color(str(item.get("color", "111111")))
    if item.get("font"):
        font.name = str(item["font"])
    elif item.get("font_face"):
        font.name = str(item["font_face"])


def compose(deck_path: Path, output_path: Path) -> None:
    deck = json.loads(deck_path.read_text(encoding="utf-8"))
    base = Path(deck.get("assets_dir") or deck_path.parent)
    if not base.is_absolute():
        base = deck_path.parent / base
    slides = deck.get("slides") or []
    if not slides:
        raise SystemExit("layered_deck.json has no slides")

    slide_w = float(deck.get("slide_width_in", 13.333333))
    slide_h = float(deck.get("slide_height_in", 7.5))
    prs = Presentation()
    prs.slide_width = Inches(slide_w)
    prs.slide_height = Inches(slide_h)
    blank_layout = prs.slide_layouts[6]

    for index, slide_spec in enumerate(slides, start=1):
        slide = prs.slides.add_slide(blank_layout)
        background = _resolve(base, slide_spec.get("background"))
        if not background or not background.exists():
            raise SystemExit(f"slide {index}: missing background image: {background}")
        slide.shapes.add_picture(str(background), 0, 0, width=prs.slide_width, height=prs.slide_height)

        frame = _resolve(base, slide_spec.get("frame"))
        if frame:
            if not frame.exists():
                raise SystemExit(f"slide {index}: missing frame image: {frame}")
            slide.shapes.add_picture(str(frame), 0, 0, width=prs.slide_width, height=prs.slide_height)

        for icon in slide_spec.get("icons") or []:
            path = _resolve(base, icon.get("file"))
            if not path or not path.exists():
                raise SystemExit(f"slide {index}: missing icon image: {path}")
            x, y, w, h = _box(icon)
            slide.shapes.add_picture(str(path), Inches(x * slide_w), Inches(y * slide_h), width=Inches(w * slide_w), height=Inches(h * slide_h))

        for item in slide_spec.get("texts") or []:
            _add_text(slide, item, slide_w, slide_h)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    print(f"saved {output_path}")


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("deck", type=Path, help="Path to layered_deck.json")
    parser.add_argument("output", type=Path, help="Output .pptx path")
    args = parser.parse_args()
    compose(args.deck, args.output)


if __name__ == "__main__":
    main()
