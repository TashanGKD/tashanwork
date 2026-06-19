#!/usr/bin/env python3
"""Compose an editable .pptx from background images, cut-out icons, and text.

Final step of Part 2 (and the merge step of Part 1). Reads one layout file
describing every slide as stacked layers (back to front):
  background (full-bleed) -> whole-frame PNG -> icon/decoration PNGs (positioned) ->
  text boxes
and writes a real PowerPoint file where text stays editable and the frame/icons
stay as movable pictures. The Image2PPTX workflow defaults to the full-slide
`frame` field. If the user explicitly asks to split the frame into movable
parts, the script also accepts `icons[]` entries with `role:"frame_part"`
generated from `frame_parts/`.

Also handles Part 1 "image deck" output: a slide with only `background` and no
icons/texts becomes a full-bleed image slide.

Usage:
    python3 scripts/compose_pptx.py deck.json out.pptx
    python3 scripts/compose_pptx.py deck.json out.pptx --preview-dir out/preview

Layout schema: see references/image-to-pptx.md. Quick form:
{
  "slide_width_in": 13.333, "slide_height_in": 7.5,
  "units": "fraction",            // "fraction" (of slide W/H) or "px"
  "ref_width": 2048, "ref_height": 1152,   // source image size; required for source_bbox/size_px QA
  "assets_dir": ".",              // base dir for relative file paths
  "slides": [
    {
      "background": "01-bg.png",
      "frame": "frame.png",
      "icons": [
        {"file": "icons/i_r1c1.png", "x":0.1,"y":0.2,"w":0.08,"h":0.08}
      ],
      "texts": [{"text":"标题","x":0.08,"y":0.05,"w":0.6,"h":0.12,
                 "size_ratio":0.074,"color":"#1A1A1A","bold":true,"align":"left","valign":"top",
                 "font":"Microsoft YaHei"}]
    }
  ]
}
A single-slide dict (without "slides") is also accepted.
"""
from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

EMU_PER_INCH = 914400


def _die(msg: str, code: int = 2):
    print(f"Error: {msg}", file=sys.stderr)
    raise SystemExit(code)


def _load_deck(path: Path):
    data = json.loads(path.read_text(encoding="utf-8"))
    if "slides" not in data:
        # Treat the whole object as a single slide; lift deck-level keys out.
        slide_keys = {"background", "frame", "shapes", "icons", "texts"}
        slide = {k: data[k] for k in slide_keys if k in data}
        deck = {k: v for k, v in data.items() if k not in slide_keys}
        deck["slides"] = [slide]
        data = deck
    data.setdefault("slide_width_in", 13.333)
    data.setdefault("slide_height_in", 7.5)
    data.setdefault("units", "fraction")
    data.setdefault("assets_dir", str(path.parent))
    return data


def _resolve(assets_dir: Path, file: str) -> Path:
    p = Path(file)
    return p if p.is_absolute() else (assets_dir / p)


def _frac(deck, item, key_xy, axis, ref):
    """Return a 0..1 fraction for x/y/w/h given the deck's unit system."""
    val = item[key_xy]
    if deck["units"] == "px":
        return val / ref
    return val


def _text_size_pt(item, sh_pt: float, ref_h: float, default=None) -> float:
    """Return text size in points.

    Preferred Image2PPTX field is size_ratio/size_pct: source text height as a
    fraction/percent of the source image height. This scales cleanly to any PPT
    slide height. Legacy size_px and absolute size(pt) remain supported.
    """
    if item.get("size") is not None:
        return float(item["size"])
    if item.get("size_ratio") is not None:
        return float(item["size_ratio"]) * sh_pt
    if item.get("size_pct") is not None:
        return float(item["size_pct"]) / 100.0 * sh_pt
    if item.get("size_px") is not None and ref_h:
        return float(item["size_px"]) * sh_pt / ref_h
    return float(default) if default is not None else 18.0


def _set_run_fonts(run, name: str):
    """Set latin + east-asian + complex-script typeface so CJK renders correctly."""
    from pptx.oxml.ns import qn
    run.font.name = name
    rpr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rpr.find(qn(tag))
        if el is None:
            el = rpr.makeelement(qn(tag), {})
            rpr.append(el)
        el.set("typeface", name)


def _set_run_alpha(run, opacity: float):
    """Set text color alpha for lightly visible GPT-vision text overlays."""
    from pptx.oxml.ns import qn
    opacity = max(0.0, min(1.0, float(opacity)))
    rpr = run._r.get_or_add_rPr()
    solid = rpr.find(qn("a:solidFill"))
    if solid is None:
        return
    srgb = solid.find(qn("a:srgbClr"))
    if srgb is None:
        return
    for old in srgb.findall(qn("a:alpha")):
        srgb.remove(old)
    alpha = srgb.makeelement(qn("a:alpha"), {"val": str(int(opacity * 100000))})
    srgb.append(alpha)


def _hex_to_rgb(value: str):
    from pptx.dml.color import RGBColor
    v = value.strip().lstrip("#")
    if len(v) == 3:
        v = "".join(ch * 2 for ch in v)
    if len(v) != 6:
        return RGBColor(0x11, 0x11, 0x11)
    return RGBColor(int(v[0:2], 16), int(v[2:4], 16), int(v[4:6], 16))


def _hex_to_tuple(value: str):
    v = value.strip().lstrip("#")
    if len(v) == 3:
        v = "".join(ch * 2 for ch in v)
    if len(v) != 6:
        return (17, 17, 17)
    return (int(v[0:2], 16), int(v[2:4], 16), int(v[4:6], 16))


def _hex_to_rgba(value: str, opacity: float = 1.0):
    r, g, b = _hex_to_tuple(value)
    return (r, g, b, int(max(0.0, min(1.0, float(opacity))) * 255))


def _set_fill_alpha(shape, opacity: float):
    """Add an <a:alpha> to a shape's solid fill so cards can be translucent."""
    from pptx.oxml.ns import qn
    spPr = shape._element.spPr
    solid = spPr.find(qn("a:solidFill"))
    if solid is None:
        return
    srgb = solid.find(qn("a:srgbClr"))
    if srgb is None:
        return
    for old in srgb.findall(qn("a:alpha")):
        srgb.remove(old)
    alpha = srgb.makeelement(qn("a:alpha"), {"val": str(int(max(0.0, min(1.0, opacity)) * 100000))})
    srgb.append(alpha)


def _add_outer_shadow(shape, blur_pt=6.0, dist_pt=3.0, alpha=0.35):
    """Best-effort soft drop shadow for card-like shapes."""
    from pptx.oxml.ns import qn
    from pptx.util import Pt
    spPr = shape._element.spPr
    for old in spPr.findall(qn("a:effectLst")):
        spPr.remove(old)
    eff = spPr.makeelement(qn("a:effectLst"), {})
    sh = eff.makeelement(qn("a:outerShdw"), {
        "blurRad": str(int(Pt(blur_pt))), "dist": str(int(Pt(dist_pt))),
        "dir": "5400000", "rotWithShape": "0"})
    clr = sh.makeelement(qn("a:srgbClr"), {"val": "000000"})
    a = clr.makeelement(qn("a:alpha"), {"val": str(int(alpha * 100000))})
    clr.append(a)
    sh.append(clr)
    eff.append(sh)
    spPr.append(eff)


def build_pptx(deck, out_path: Path):
    from pptx import Presentation
    from pptx.util import Emu, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
    from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

    shape_map = {"rect": MSO_SHAPE.RECTANGLE, "rounded_rect": MSO_SHAPE.ROUNDED_RECTANGLE,
                 "oval": MSO_SHAPE.OVAL, "ellipse": MSO_SHAPE.OVAL,
                 "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE, "chevron": MSO_SHAPE.CHEVRON,
                 "right_arrow": MSO_SHAPE.RIGHT_ARROW, "left_arrow": MSO_SHAPE.LEFT_ARROW,
                 "up_arrow": MSO_SHAPE.UP_ARROW, "down_arrow": MSO_SHAPE.DOWN_ARROW,
                 "pentagon": MSO_SHAPE.REGULAR_PENTAGON, "hexagon": MSO_SHAPE.HEXAGON,
                 "parallelogram": MSO_SHAPE.PARALLELOGRAM, "trapezoid": MSO_SHAPE.TRAPEZOID,
                 "diamond": MSO_SHAPE.DIAMOND}

    assets_dir = Path(deck["assets_dir"])
    sw_in = float(deck["slide_width_in"])
    sh_in = float(deck["slide_height_in"])
    sw_emu = int(round(sw_in * EMU_PER_INCH))
    sh_emu = int(round(sh_in * EMU_PER_INCH))
    ref_w = float(deck.get("ref_width") or 0)
    ref_h = float(deck.get("ref_height") or 0)
    sh_pt = sh_in * 72.0

    prs = Presentation()
    prs.slide_width = Emu(sw_emu)
    prs.slide_height = Emu(sh_emu)
    blank = prs.slide_layouts[6]

    align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                 "right": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY}
    anchor_map = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE,
                  "center": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}

    for idx, sl in enumerate(deck["slides"], 1):
        slide = prs.slides.add_slide(blank)

        bg = sl.get("background")
        if bg:
            bg_path = _resolve(assets_dir, bg)
            if not bg_path.exists():
                _die(f"slide {idx}: background not found: {bg_path}")
            slide.shapes.add_picture(str(bg_path), 0, 0, width=Emu(sw_emu), height=Emu(sh_emu))

        frame = sl.get("frame")
        if frame:
            fr_path = _resolve(assets_dir, frame)
            if not fr_path.exists():
                _die(f"slide {idx}: frame not found: {fr_path}")
            slide.shapes.add_picture(str(fr_path), 0, 0, width=Emu(sw_emu), height=Emu(sh_emu))

        for shp in sl.get("shapes", []):
            fx = _frac(deck, shp, "x", "x", ref_w)
            fy = _frac(deck, shp, "y", "y", ref_h)
            fw = _frac(deck, shp, "w", "w", ref_w)
            fh = _frac(deck, shp, "h", "h", ref_h)
            left, top = Emu(int(fx * sw_emu)), Emu(int(fy * sh_emu))
            wid, hei = Emu(int(fw * sw_emu)), Emu(int(fh * sh_emu))
            stype = shp.get("type", "rounded_rect")

            if stype == "line":
                conn = slide.shapes.add_connector(
                    MSO_CONNECTOR.STRAIGHT, left, top, Emu(int((fx + fw) * sw_emu)),
                    Emu(int((fy + fh) * sh_emu)))
                conn.line.color.rgb = _hex_to_rgb(shp.get("line", shp.get("fill", "#FFFFFF")))
                conn.line.width = Pt(float(shp.get("line_width", 1.5)))
                continue

            shape = slide.shapes.add_shape(shape_map.get(stype, MSO_SHAPE.ROUNDED_RECTANGLE),
                                           left, top, wid, hei)
            if stype == "rounded_rect" and "radius" in shp:
                try:
                    shape.adjustments[0] = float(shp["radius"])
                except Exception:
                    pass

            if shp.get("fill"):
                shape.fill.solid()
                shape.fill.fore_color.rgb = _hex_to_rgb(shp["fill"])
                if shp.get("opacity") is not None:
                    _set_fill_alpha(shape, float(shp["opacity"]))
            else:
                shape.fill.background()

            if shp.get("line"):
                shape.line.color.rgb = _hex_to_rgb(shp["line"])
                shape.line.width = Pt(float(shp.get("line_width", 1.0)))
            else:
                shape.line.fill.background()

            if shp.get("shadow"):
                _add_outer_shadow(shape)
            else:
                shape.shadow.inherit = False

            if shp.get("rotation"):
                shape.rotation = float(shp["rotation"])

        for ic in sl.get("icons", []):
            ip = _resolve(assets_dir, ic["file"])
            if not ip.exists():
                _die(f"slide {idx}: icon not found: {ip}")
            fx = _frac(deck, ic, "x", "x", ref_w)
            fy = _frac(deck, ic, "y", "y", ref_h)
            fw = _frac(deck, ic, "w", "w", ref_w)
            fh = _frac(deck, ic, "h", "h", ref_h)
            slide.shapes.add_picture(
                str(ip), Emu(int(fx * sw_emu)), Emu(int(fy * sh_emu)),
                width=Emu(int(fw * sw_emu)), height=Emu(int(fh * sh_emu)))

        for tx in sl.get("texts", []):
            fx = _frac(deck, tx, "x", "x", ref_w)
            fy = _frac(deck, tx, "y", "y", ref_h)
            fw = _frac(deck, tx, "w", "w", ref_w)
            fh = _frac(deck, tx, "h", "h", ref_h)
            box = slide.shapes.add_textbox(
                Emu(int(fx * sw_emu)), Emu(int(fy * sh_emu)),
                Emu(int(fw * sw_emu)), Emu(int(fh * sh_emu)))
            tf = box.text_frame
            tf.word_wrap = True
            if tx.get("fit_text", True):
                tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            tf.vertical_anchor = anchor_map.get(tx.get("valign", "top"), MSO_ANCHOR.TOP)
            # PowerPoint text boxes have non-zero default internal margins. The
            # layout bbox already represents the visual text box measured on the
            # source image, so default margins create systematic placement drift.
            tf.margin_left = Pt(float(tx.get("margin_left", 0)))
            tf.margin_right = Pt(float(tx.get("margin_right", 0)))
            tf.margin_top = Pt(float(tx.get("margin_top", 0)))
            tf.margin_bottom = Pt(float(tx.get("margin_bottom", 0)))

            size_pt = _text_size_pt(tx, sh_pt, ref_h)

            color = _hex_to_rgb(tx.get("color", "#111111"))
            font = tx.get("font", "Microsoft YaHei")
            bold = bool(tx.get("bold", False))
            italic = bool(tx.get("italic", False))
            align = align_map.get(tx.get("align", "left"), PP_ALIGN.LEFT)
            line_spacing = tx.get("line_spacing")
            opacity = float(tx.get("opacity", 1.0))
            if tx.get("name"):
                box.name = str(tx["name"])

            runs = tx.get("runs")
            if runs:
                para = tf.paragraphs[0]
                para.alignment = align
                if line_spacing:
                    para.line_spacing = float(line_spacing)
                for rinfo in runs:
                    run = para.add_run()
                    run.text = str(rinfo.get("text", ""))
                    run.font.size = Pt(_text_size_pt(rinfo, sh_pt, ref_h, default=size_pt))
                    run.font.bold = bool(rinfo.get("bold", bold))
                    run.font.italic = bool(rinfo.get("italic", italic))
                    run.font.color.rgb = _hex_to_rgb(rinfo["color"]) if rinfo.get("color") else color
                    _set_run_fonts(run, str(rinfo.get("font", font)))
                    _set_run_alpha(run, float(rinfo.get("opacity", opacity)))
            else:
                lines = str(tx.get("text", "")).split("\n")
                for li, line in enumerate(lines):
                    para = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
                    para.alignment = align
                    if line_spacing:
                        para.line_spacing = float(line_spacing)
                    run = para.add_run()
                    run.text = line
                    run.font.size = Pt(size_pt)
                    run.font.bold = bold
                    run.font.italic = italic
                    run.font.color.rgb = color
                    _set_run_fonts(run, font)
                    _set_run_alpha(run, opacity)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    print(f"Wrote {out_path}  ({len(deck['slides'])} slides)")


# --------------------------- optional PNG preview ---------------------------

def _find_cjk_font(bold=False):
    candidates_bold = [
        "C:/Windows/Fonts/msyhbd.ttc",
        "C:/Windows/Fonts/simhei.ttf",
        "C:/Windows/Fonts/simsun.ttc",
        "/System/Library/Fonts/PingFang.ttc",
        "/System/Library/Fonts/STHeiti Medium.ttc",
        "/usr/share/fonts/opentype/noto/NotoSansCJK-Bold.ttc",
        "/usr/share/fonts/truetype/noto/NotoSansCJK-Bold.ttc",
    ]
    candidates = [
        "C:/Windows/Fonts/msyh.ttc",
        "C:/Windows/Fonts/msyhbd.ttc",
        "C:/Windows/Fonts/simhei.ttf",
        "C:/Windows/Fonts/simsun.ttc",
        "C:/Windows/Fonts/simkai.ttf",
        "/System/Library/Fonts/PingFang.ttc",
        "/System/Library/Fonts/STHeiti Light.ttc",
        "/Library/Fonts/Arial Unicode.ttf",
        "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/truetype/noto/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    ]
    pool = candidates_bold + candidates if bold else candidates
    for c in pool:
        if Path(c).exists():
            return c
    return None


def _poly_points(stype, x0, y0, x1, y1):
    """Approximate polygon vertices for preview rendering of non-rect shapes."""
    w, h = x1 - x0, y1 - y0
    if stype == "triangle":
        return [(x0 + w / 2, y0), (x1, y1), (x0, y1)]
    if stype == "diamond":
        return [(x0 + w / 2, y0), (x1, y0 + h / 2), (x0 + w / 2, y1), (x0, y0 + h / 2)]
    if stype == "right_arrow":
        m = h * 0.30
        return [(x0, y0 + m), (x1 - w * 0.4, y0 + m), (x1 - w * 0.4, y0), (x1, y0 + h / 2),
                (x1 - w * 0.4, y1), (x1 - w * 0.4, y1 - m), (x0, y1 - m)]
    if stype == "left_arrow":
        m = h * 0.30
        return [(x1, y0 + m), (x0 + w * 0.4, y0 + m), (x0 + w * 0.4, y0), (x0, y0 + h / 2),
                (x0 + w * 0.4, y1), (x0 + w * 0.4, y1 - m), (x1, y1 - m)]
    if stype == "up_arrow":
        m = w * 0.30
        return [(x0 + m, y1), (x0 + m, y0 + h * 0.4), (x0, y0 + h * 0.4), (x0 + w / 2, y0),
                (x1, y0 + h * 0.4), (x1 - m, y0 + h * 0.4), (x1 - m, y1)]
    if stype == "chevron":
        notch = w * 0.25
        return [(x0, y0), (x1 - notch, y0), (x1, y0 + h / 2), (x1 - notch, y1),
                (x0, y1), (x0 + notch, y0 + h / 2)]
    if stype == "trapezoid":
        inset = w * 0.22
        return [(x0 + inset, y0), (x1 - inset, y0), (x1, y1), (x0, y1)]
    if stype == "parallelogram":
        sk = w * 0.22
        return [(x0 + sk, y0), (x1, y0), (x1 - sk, y1), (x0, y1)]
    if stype == "pentagon":
        import math
        cx, cy, rx, ry = x0 + w / 2, y0 + h / 2, w / 2, h / 2
        return [(cx + rx * math.sin(2 * math.pi * i / 5),
                 cy - ry * math.cos(2 * math.pi * i / 5)) for i in range(5)]
    if stype == "hexagon":
        return [(x0 + w * 0.25, y0), (x0 + w * 0.75, y0), (x1, y0 + h / 2),
                (x0 + w * 0.75, y1), (x0 + w * 0.25, y1), (x0, y0 + h / 2)]
    return None


def _wrap_text(draw, text, font, max_w):
    trailing_punctuation = set("，。；：、！？,.!?;:)]}）】》")
    out = []
    for raw in text.split("\n"):
        if not raw:
            out.append("")
            continue
        line = ""
        for ch in raw:
            trial = line + ch
            w = draw.textlength(trial, font=font)
            if w > max_w and line:
                if ch in trailing_punctuation:
                    line = trial
                else:
                    out.append(line)
                    line = ch
            else:
                line = trial
        out.append(line)
    return out


def render_previews(deck, preview_dir: Path):
    try:
        from PIL import Image, ImageDraw, ImageFont
    except ImportError:
        print("Preview skipped: Pillow not available.", file=sys.stderr)
        return

    assets_dir = Path(deck["assets_dir"])
    sw_in = float(deck["slide_width_in"])
    sh_in = float(deck["slide_height_in"])
    ratio = sw_in / sh_in
    if deck["units"] == "px" and deck.get("ref_width") and deck.get("ref_height"):
        CW, CH = int(deck["ref_width"]), int(deck["ref_height"])
    else:
        CW = 1600
        CH = int(round(CW / ratio))
    sh_pt = sh_in * 72.0
    ref_w = float(deck.get("ref_width") or CW)
    ref_h = float(deck.get("ref_height") or CH)

    preview_dir.mkdir(parents=True, exist_ok=True)
    regular_path = _find_cjk_font(bold=False)
    bold_path = _find_cjk_font(bold=True)

    for idx, sl in enumerate(deck["slides"], 1):
        canvas = Image.new("RGBA", (CW, CH), (255, 255, 255, 255))
        bg = sl.get("background")
        if bg:
            bp = _resolve(assets_dir, bg)
            if bp.exists():
                with Image.open(bp) as im:
                    canvas.paste(im.convert("RGBA").resize((CW, CH)), (0, 0))

        frame = sl.get("frame")
        if frame:
            fp = _resolve(assets_dir, frame)
            if fp.exists():
                with Image.open(fp) as im:
                    canvas.alpha_composite(im.convert("RGBA").resize((CW, CH)), (0, 0))

        for shp in sl.get("shapes", []):
            fx = _frac(deck, shp, "x", "x", ref_w)
            fy = _frac(deck, shp, "y", "y", ref_h)
            fw = _frac(deck, shp, "w", "w", ref_w)
            fh = _frac(deck, shp, "h", "h", ref_h)
            x0, y0 = int(fx * CW), int(fy * CH)
            x1, y1 = int((fx + fw) * CW), int((fy + fh) * CH)
            overlay = Image.new("RGBA", (CW, CH), (0, 0, 0, 0))
            od = ImageDraw.Draw(overlay)
            stype = shp.get("type", "rounded_rect")
            a = int(float(shp.get("opacity", 1.0)) * 255) if shp.get("fill") else 0
            fill = _hex_to_tuple(shp["fill"]) + (a,) if shp.get("fill") else None
            line = _hex_to_tuple(shp["line"]) + (255,) if shp.get("line") else None
            lw = int(round(float(shp.get("line_width", 1.0)) * CH / (sh_in * 72.0)))
            lw = max(1, lw) if line else 0
            poly = _poly_points(stype, x0, y0, x1, y1)
            if stype == "line":
                od.line([x0, y0, x1, y1], fill=line or (255, 255, 255, 255),
                        width=max(1, lw or 2))
            elif stype in ("oval", "ellipse"):
                od.ellipse([x0, y0, x1, y1], fill=fill, outline=line, width=lw)
            elif stype == "rounded_rect":
                rad = int(float(shp.get("radius", 0.12)) * min(x1 - x0, y1 - y0))
                od.rounded_rectangle([x0, y0, x1, y1], radius=max(1, rad), fill=fill,
                                     outline=line, width=lw)
            elif poly is not None:
                od.polygon(poly, fill=fill, outline=line, width=lw)
            else:
                od.rectangle([x0, y0, x1, y1], fill=fill, outline=line, width=lw)
            canvas.alpha_composite(overlay)

        for ic in sl.get("icons", []):
            ip = _resolve(assets_dir, ic["file"])
            if not ip.exists():
                continue
            fx = _frac(deck, ic, "x", "x", ref_w)
            fy = _frac(deck, ic, "y", "y", ref_h)
            fw = _frac(deck, ic, "w", "w", ref_w)
            fh = _frac(deck, ic, "h", "h", ref_h)
            tw, th = max(1, int(fw * CW)), max(1, int(fh * CH))
            with Image.open(ip) as im:
                icon = im.convert("RGBA").resize((tw, th))
            canvas.alpha_composite(icon, (int(fx * CW), int(fy * CH)))

        draw = ImageDraw.Draw(canvas)
        for tx in sl.get("texts", []):
            fx = _frac(deck, tx, "x", "x", ref_w)
            fy = _frac(deck, tx, "y", "y", ref_h)
            fw = _frac(deck, tx, "w", "w", ref_w)
            size_pt = _text_size_pt(tx, sh_pt, ref_h)
            px = max(8, int(round(size_pt * CH / sh_pt)))
            bold = bool(tx.get("bold", False))
            fpath = (bold_path if bold else regular_path) or regular_path
            try:
                font = ImageFont.truetype(fpath, px) if fpath else ImageFont.load_default()
            except Exception:
                font = ImageFont.load_default()
            color = tx.get("color", "#111111")
            align = tx.get("align", "left")
            valign = tx.get("valign", "top")
            opacity = float(tx.get("opacity", 1.0))
            bx0, by0 = int(fx * CW), int(fy * CH)
            bw, bh = int(fw * CW), int(_frac(deck, tx, "h", "h", ref_h) * CH)
            line_h = int(px * 1.3)
            runs = tx.get("runs")
            if runs:
                segs = []
                total = 0.0
                for rinfo in runs:
                    rb = bool(rinfo.get("bold", bold))
                    rpx = max(8, int(round(_text_size_pt(rinfo, sh_pt, ref_h, default=size_pt) * CH / sh_pt)))
                    rfp = (bold_path if rb else regular_path) or regular_path
                    try:
                        rfont = ImageFont.truetype(rfp, rpx) if rfp else ImageFont.load_default()
                    except Exception:
                        rfont = ImageFont.load_default()
                    rtext = str(rinfo.get("text", ""))
                    w = draw.textlength(rtext, font=rfont)
                    segs.append((rtext, rfont, rinfo.get("color", color),
                                 float(rinfo.get("opacity", opacity)), rb, w))
                    total += w
                cx = bx0 + (bw - total if align == "right" else
                            (bw - total) // 2 if align == "center" else 0)
                cx = max(bx0, cx)
                ty = by0 + (max(0, bh - line_h) if valign == "bottom" else
                            max(0, (bh - line_h) // 2) if valign in ("middle", "center") else 0)
                for rtext, rfont, rcolor, ropacity, rb, w in segs:
                    stroke = 1 if rb and not bold_path else 0
                    fill = _hex_to_rgba(rcolor, ropacity)
                    draw.text((cx, ty), rtext, fill=fill, font=rfont,
                              stroke_width=stroke, stroke_fill=fill)
                    cx += w
            else:
                text_value = str(tx.get("text", ""))
                if tx.get("fit_text", True):
                    while px > 6:
                        lines = _wrap_text(draw, text_value, font, max(1, bw))
                        line_h = int(px * 1.3)
                        total_h = line_h * max(1, len(lines))
                        max_line_w = max((draw.textlength(line, font=font) for line in lines), default=0)
                        if total_h <= max(1, bh) and max_line_w <= max(1, bw):
                            break
                        px -= 1
                        try:
                            font = ImageFont.truetype(fpath, px) if fpath else ImageFont.load_default()
                        except Exception:
                            font = ImageFont.load_default()
                    line_h = int(px * 1.3)
                lines = _wrap_text(draw, text_value, font, max(1, bw))
                total_h = line_h * max(1, len(lines))
                ty = by0 + (max(0, bh - total_h) if valign == "bottom" else
                            max(0, (bh - total_h) // 2) if valign in ("middle", "center") else 0)
                stroke = 1 if bold and not bold_path else 0
                for li, line in enumerate(lines):
                    lw_ = draw.textlength(line, font=font)
                    lx = bx0 + (bw - lw_ if align == "right" else
                                (bw - lw_) / 2 if align == "center" else 0)
                    fill = _hex_to_rgba(color, opacity)
                    draw.text((max(bx0, int(lx)), ty + li * line_h), line, fill=fill,
                              font=font, stroke_width=stroke, stroke_fill=fill)

        out = preview_dir / f"slide_{idx:02d}.png"
        canvas.convert("RGB").save(out)
        print(f"Preview: {out}")


def main() -> None:
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("layout", help="deck.json / layout.json path.")
    ap.add_argument("out", help="Output .pptx path.")
    ap.add_argument("--preview-dir", help="If set, also render PNG previews here for QA.")
    args = ap.parse_args()

    lp = Path(args.layout)
    if not lp.exists():
        _die(f"layout file not found: {lp}")
    deck = _load_deck(lp)
    build_pptx(deck, Path(args.out))
    if args.preview_dir:
        render_previews(deck, Path(args.preview_dir))


if __name__ == "__main__":
    main()
